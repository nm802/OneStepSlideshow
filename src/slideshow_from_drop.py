from __future__ import annotations
import os
import sys
import datetime
import typing

from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.slide import Slides


class Rectangle:
    def __init__(self, width, height):
        self.width = int(width)
        self.height = int(height)
        self.aspect_ratio = width / height

    def fit(self, target_shape: Rectangle) -> None:
        """
        Reset width and height to fit the target rectangle.
        Args:
            target_shape: Target Rectangle class instance

        Returns: Rectangle class instance after resize.
        """
        if self.aspect_ratio > target_shape.aspect_ratio:  # 横長
            self.width = target_shape.width
            self.height = self.width / self.aspect_ratio
        else:  # 縦長
            self.height = target_shape.height
            self.width = self.aspect_ratio * self.height

    def fill(self, target_shape: Rectangle) -> None:
        """
        Reset width and height to fill the target rectangle.
        Args:
            target_shape: Target Rectangle class instance

        Returns: Rectangle class instance after resize.

        """
        if self.aspect_ratio > target_shape.aspect_ratio:  # 横長
            self.height = target_shape.height
            self.width = self.aspect_ratio * self.height
        else:  # 縦長
            self.width = target_shape.width
            self.height = self.width / self.aspect_ratio


def add_slide(prs: Presentation()) -> Slides:
    # 白紙スライドの追加(ID=6は白紙スライド)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide


def add_picture_fill(img_path: str, slide, left, top, grid_shape: Rectangle, img_shape: Rectangle) -> typing.NoReturn:
    """
    画像を(中央crop後に指定位置にくるよう)addした後，gridサイズにcropする
    """
    crop = {'H': 0, 'V': 0}
    # fillモードの場合，画像のほうが横または縦のどちらかがgridより大きい
    if img_shape.width > grid_shape.width:
        trim_width = (img_shape.width - grid_shape.width) / 2
        left = left + trim_width
        crop['H'] = trim_width / img_shape.width
    elif img_shape.height > grid_shape.height:
        trim_width = (img_shape.height - grid_shape.height) / 2
        top = top + trim_width
        crop['V'] = trim_width / img_shape.height

    # add_shape
    shape = slide.shapes.add_picture(img_path, left, top, width=img_shape.width)

    # crop
    shape.crop_top = crop['V']
    shape.crop_bottom = crop['V']
    shape.crop_left = crop['H']
    shape.crop_right = crop['H']
    shape.width = grid_shape.width
    shape.height = grid_shape.height


def add_picture_fit(img_path: str, slide, left, top, grid_shape: Rectangle, img_shape: Rectangle) -> typing.NoReturn:
    """
    画像を(中央crop後に指定位置にくるよう)addした後，gridサイズにcropする
    """
    add = {'H': 0, 'V': 0}
    if img_shape.width < grid_shape.width:
        add_width = (grid_shape.width - img_shape.width) / 2
        left = left - add_width
        add['H'] = - add_width / img_shape.width
    else:
        add_width = (grid_shape.height - img_shape.height) / 2
        top = top - add_width
        add['V'] = - add_width / img_shape.height

    # add_shape
    shape = slide.shapes.add_picture(img_path, left, top, width=img_shape.width)

    shape.crop_top = add['V']
    shape.crop_bottom = add['V']
    shape.crop_left = add['H']
    shape.crop_right = add['H']
    shape.width = grid_shape.width
    shape.height = grid_shape.height


def add_filename(file_name: str, slide: Slides, left, top, grid_shape: Rectangle):
    """
    ファイル名を各写真右上に貼り付ける
    """

    text_box = slide.shapes.add_textbox(left, top, grid_shape.width, Pt(16))
    text_frame = text_box.text_frame
    text_frame.text = file_name
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    return slide


def make_slideshow(_img_file_paths: list, _slide_aspect_ratio: float = 4 / 3, _grid_definition: tuple = (1, 1),
                   _mode: str = 'fill', _with_filename: bool = True):
    """
    4:3 (default) 9144000x6858000, 16:9 9144000x5143680
    Unit: English Metric Units = 1/360000 of centimeter
    Width is always 25.4 cm = 25.4 x 360000 EMU = 9144000 EMU.
    Args:
        _img_file_paths: list of image file paths (full path)
        _slide_aspect_ratio: slide width/height
        _grid_definition: image grid numbers as tuple (row, column)
        _mode: 'fill' or 'fit'. images cropped in fill mode.
        _with_filename: True to add filename on each image.

    Returns:
        No returns.

    """

    # region Check Input
    if (len(_img_file_paths) == 0) or (type(_img_file_paths) is not list):
        print('no image file included. valid extensions: png/jpg/jpeg/bmp only.')
        return
    if (_mode != 'fit') and (_mode != 'fill'):
        print('invalid _mode thrown. _mode must be "fill" or "fit".')
        return
    # endregion

    # region Initialize Grid
    slide_shape = Rectangle(9144000, 9144000 / _slide_aspect_ratio)
    # Generate presentation object
    prs = Presentation()
    prs.slide_width = int(slide_shape.width)
    prs.slide_height = int(slide_shape.height)
    # Calculate grid size/number
    grid_shape = Rectangle(slide_shape.width / _grid_definition[1], slide_shape.height / _grid_definition[0])
    qty_per_a_slide = int(_grid_definition[0]) * int(_grid_definition[1])
    # endregion

    for i, img_path in enumerate(_img_file_paths):
        im = Image.open(img_path)
        im_width, im_height = im.size
        img_shape = Rectangle(im_width, im_height)
        # img_shape Rectangleオブジェクトのfitまたはfillメソッドの実行
        getattr(img_shape, _mode)(grid_shape)

        if i % qty_per_a_slide == 0:
            slide = add_slide(prs)
        position = i % qty_per_a_slide

        # グリッドの左上座標
        grid_left = grid_shape.width * (position % _grid_definition[1])
        grid_top = grid_shape.height * int(position / _grid_definition[1])
        left = grid_left + ((grid_shape.width - img_shape.width) / 2)
        top = grid_top + ((grid_shape.height - img_shape.height) / 2)

        # 画像をスライドに追加
        if _mode == 'fit':
            add_picture_fit(img_path, slide, left, top, grid_shape, img_shape)
        else:
            add_picture_fill(img_path, slide, left, top, grid_shape, img_shape)

        # ファイル名をスライドに追加
        if _with_filename:
            add_filename(os.path.basename(img_path), slide, grid_left, grid_top, grid_shape)

    # save .pptx file
    output_dir = os.path.dirname(_img_file_paths[0])
    print('Output .pptx file dir = ' + output_dir)
    output_file_name = output_dir + '/slideshow_' + datetime.datetime.now().strftime('%Y%m%d_%H%M%S') + '.pptx'
    prs.save(output_file_name)


def filepath_list(_paths: list) -> list:
    """
    与えられたファイルパスのリストから画像ファイルだけを返す。
    ディレクトリが含まれていたらその中のファイルを再帰的に取得する。

    Args:
        _paths: ファイルパスのリスト。batファイルから与えられることを想定

    Returns:
        画像ファイルパスのリスト
    """
    # return
    _img_file_paths = list()

    # filter
    for p in _paths:
        if os.path.exists(p) and p.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp')):
            _img_file_paths.append(p)
        elif os.path.isdir(p):
            _img_file_paths += filepath_list([os.path.join(p, f) for f in os.listdir(p)])

    return _img_file_paths


if __name__ == '__main__':
    # コマンドライン引数として以下渡す
    # args[0]: slideshow_from_drop.py (このファイルの名前)
    # args[1]: 0 or 1; 0 -> _slide_aspect_ratio = 4 / 3, 1 -> _slide_aspect_ratio = 16 / 9
    # args[2]: row number of grid. Prompted if args[2] == 0
    # args[3]: column number of grid. Prompted if args[2] == 0
    # args[4]: 0 or 1; 0 -> _mode = 'fill', 1 -> _mode = 'fit'
    # args[5]: y or n; y -> _with_filename = True, n -> _with_filename = False. Prompted if args[5] == 0
    # args[6]~: target file paths

    args = sys.argv
    if len(args) < 6:
        print('Too short args. Confirm args definition.')
        sys.exit()

    if int(args[1]) == 0:
        slide_aspect_ratio = 4 / 3
    elif int(args[1]) == 1:
        slide_aspect_ratio = 16 / 9
    else:
        print('Arg 1 is wrong value. Value = ' + args[1])
        sys.exit()

    if int(args[2]) == 0:
        r = input('Enter Grid Row Number: ')
        c = input('Enter Grid Column Number: ')
        grid_definition = ((int(r), int(c)))
    else:
        grid_definition = (int(args[2]), int(args[3]))

    if int(args[4]) == 0:
        mode = 'fill'
    elif int(args[4]) == 1:
        mode = 'fit'
    else:
        print('Arg 4 is wrong value. Value = ' + args[4])
        sys.exit()

    answer = input('With Filename? [y/n]: ') if int(args[5]) == 0 else args[5]
    if str(answer).lower() == 'n':
        with_filename = False
    else:
        with_filename = True

    img_file_paths = filepath_list(args[6:])

    make_slideshow(img_file_paths, _slide_aspect_ratio=slide_aspect_ratio, _grid_definition=grid_definition, _mode=mode,
                   _with_filename=with_filename)
