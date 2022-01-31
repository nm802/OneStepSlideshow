from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from PIL import Image
import os
import datetime

# スライドサイズ
SLIDE_WIDTH, SLIDE_HEIGHT = 9144000, 6858000  # 4:3 (default) 9144000x6858000
#SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000 #  16:9
# スライド中心のX、Y座標（左上が原点）
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH / 2, SLIDE_HEIGHT / 2
# スライドのアスペクト比
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT

# 画像の格納ディレクトリ
target_img_dir = 'C:/Users/s060054/Desktop/220128_BottomCaseフック/microscope/'

# 出力ファイル名
OUTPUT_FILE_PATH = target_img_dir + 'slideshow_' + datetime.datetime.now().strftime('%Y%m%d_%H%M%S') + '.pptx'


# 受け取ったプレゼンテーションオブジェクトにスライドを追加し、追加されたスライドオブジェクトを返す。
def add_slide(prs):
    # 白紙スライドの追加(ID=6は白紙スライド)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide


# 画像をスライド中心に貼り付ける
def add_picture(slide, img_file):
    # 画像サイズを取得してアスペクト比を得る
    im = Image.open(img_file)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    # スライドと画像のアスペクト比に応じて処理を分岐
    # 画像のほうが横長だったら横めいっぱいに広げる
    if aspect_ratio > SLIDE_ASPECT_RATIO:
        img_display_width = SLIDE_WIDTH
        img_display_height = img_display_width / aspect_ratio
    else:  # 像のほうが縦長だったら縦めいっぱいに広げる
        img_display_height = SLIDE_HEIGHT
        img_display_width = img_display_height * aspect_ratio
    # センタリングする場合の画像の左上座標を計算
    left = IMG_CENTER_X - img_display_width / 2
    top = IMG_CENTER_Y - img_display_height / 2

    # 画像をスライドに追加
    if aspect_ratio > SLIDE_ASPECT_RATIO:
        slide.shapes.add_picture(img_file, left, top, width=img_display_width)
    else:
        slide.shapes.add_picture(img_file, left, top, height=img_display_height)

    return slide


# ファイル名をスライド右上に貼り付ける
def add_filename(slide, fname):
    text_box = slide.shapes.add_textbox(0, 0, SLIDE_WIDTH, Pt(28))
    # text_box.fill.solid()
    # text_box.fill.fore_color.rgb = RGBColor(0, 0, 0)
    text_frame = text_box.text_frame
    text_frame.text = fname
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    return slide


# スライドオブジェクトの定義
prs = Presentation()
# スライドサイズの指定
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# img画像のファイル名を取得
img_files = os.listdir(target_img_dir)
# pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
img_files = [name for name in img_files if (name.endswith(".png") or name.endswith(".jpg") or name.endswith(".jpeg")
                                            or name.endswith(".bmp"))]
# 昇順にソート（この順番でスライドに貼り付けられる）
img_files.sort()  # 昇順にsort

for name in img_files:
    path = target_img_dir + name
    slide = add_slide(prs)
    add_picture(slide, path)
    add_filename(slide, name)
# pptxファイルを出力する
prs.save(OUTPUT_FILE_PATH)
